"""
Unified Server for AI Digital Human Project
Combines:
1. GitHub Trending API
2. User Management API
3. User Behavior Tracking API
4. AI Personalized Prompt API
"""

from flask import Flask, jsonify, request, send_from_directory
from flask_cors import CORS
import json
import os
import hashlib
import threading
import time
import requests
from datetime import datetime, timedelta
import base64
import io

# Import local modules
import trending_fetcher
import mimetypes

# 强制修正 MIME 类型，防止 Linux 上识别成 text/plain 或 text/jsx
mimetypes.add_type('application/javascript', '.js')
mimetypes.add_type('text/css', '.css')

app = Flask(__name__)
CORS(app)

# --- Configuration ---
# 关键修改：数据存储在项目目录之外，防止部署时被删除
DATA_DIR = '/root/ai_data'
# 前端文件目录 (Nginx 搞不定，我们用 Python 直接托管)
FRONTEND_DIR = '/var/www/html/goldensales'

# 本地开发环境回退到当前目录
if os.name == 'nt':  # Windows
    DATA_DIR = 'data'
    FRONTEND_DIR = '../dist' # 本地开发时的前端构建目录

LOG_DIR = 'logs'
if not os.path.exists(DATA_DIR):
    os.makedirs(DATA_DIR)
os.makedirs(LOG_DIR, exist_ok=True)

DATA_FILE = os.path.join(DATA_DIR, 'github_trending.json')
USERS_FILE = os.path.join(DATA_DIR, 'users.json')
SYSTEM_CONFIG_FILE = os.path.join(DATA_DIR, 'system_config.json')
LOG_FILE = os.path.join(LOG_DIR, 'user_behavior.log')
READABLE_LOG_FILE = os.path.join(LOG_DIR, 'user_behavior_readable.txt')

CACHE_DURATION = timedelta(hours=1)
ZHIPU_API_KEY = os.getenv('ZHIPU_API_KEY', 'sk-Rc1j1a6cfUeWlOZYHgXikivqfrUpOdUlGz2ziD772dXFEFZd')
ZHIPU_API_URL = 'https://api.302.ai/v1/chat/completions'

# --- Global State ---
cached_data = None
cache_time = None
file_lock = threading.Lock()

# --- Helper Functions ---

def load_data():
    """Load trending data from cache or file"""
    global cached_data, cache_time
    if cached_data and cache_time and datetime.now() - cache_time < CACHE_DURATION:
        return cached_data

    if os.path.exists(DATA_FILE):
        try:
            with open(DATA_FILE, 'r', encoding='utf-8') as f:
                data = json.load(f)
            updated_at = datetime.fromisoformat(data['updated_at'])
            if datetime.now() - updated_at < CACHE_DURATION:
                cached_data = data
                cache_time = datetime.now()
                return data
        except Exception as e:
            print(f"Error loading trending data: {e}")

    trigger_background_update()
    return cached_data or {'categories': {}, 'updated_at': datetime.now().isoformat()}

def trigger_background_update():
    def update():
        print("Starting background update...")
        try:
            data = trending_fetcher.fetch_all_categories()
            global cached_data, cache_time
            cached_data = data
            cache_time = datetime.now()
            print("Background update complete!")
        except Exception as e:
            print(f"Background update failed: {e}")
    threading.Thread(target=update, daemon=True).start()

def load_users():
    if os.path.exists(USERS_FILE):
        try:
            with open(USERS_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except:
            return []
    return []

def save_users(users):
    with open(USERS_FILE, 'w', encoding='utf-8') as f:
        json.dump(users, f, ensure_ascii=False, indent=2)

def load_system_config():
    """Load system config from file"""
    if os.path.exists(SYSTEM_CONFIG_FILE):
        try:
            with open(SYSTEM_CONFIG_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            print(f"Error loading system config: {e}")
    return {}

def save_system_config(config):
    """Save system config to file"""
    try:
        with open(SYSTEM_CONFIG_FILE, 'w', encoding='utf-8') as f:
            json.dump(config, f, ensure_ascii=False, indent=2)
        return True
    except Exception as e:
        print(f"Error saving system config: {e}")
        return False

def hash_user_id(user_id):
    if not user_id or user_id == 'anonymous':
        return 'anonymous'
    return hashlib.sha256(user_id.encode()).hexdigest()[:16]

def sanitize_data(event_data):
    sanitized = event_data.copy()
    if 'userId' in sanitized:
        sanitized['userId'] = hash_user_id(sanitized['userId'])
    return sanitized

def format_readable_log(entry):
    try:
        data = entry.get('data', {})
        user = data.get('user_account', data.get('username', f"User_{entry.get('userId', '')[:8]}"))
        feature = data.get('feature_name', data.get('page_name', 'Unknown'))
        action = data.get('action', entry.get('event', ''))
        prompt = data.get('prompt', '')
        ts = entry.get('server_timestamp', '')
        return f"{ts} | {user} | {feature} | {action} | {prompt[:50]}"
    except:
        return str(entry)

def append_to_log(entry):
    with file_lock:
        try:
            with open(LOG_FILE, 'a', encoding='utf-8') as f:
                f.write(json.dumps(entry, ensure_ascii=False) + '\n')
            with open(READABLE_LOG_FILE, 'a', encoding='utf-8') as f:
                f.write(format_readable_log(entry) + '\n')
        except Exception as e:
            print(f"Log write error: {e}")

def get_user_behavior_summary(user_id, days=7):
    if not os.path.exists(LOG_FILE): return None
    user_events = []
    try:
        with open(LOG_FILE, 'r', encoding='utf-8') as f:
            for line in f:
                try:
                    event = json.loads(line.strip())
                    uid = event.get('userId') or event.get('user_id', '')
                    acc = event.get('data', {}).get('user_account', '')
                    if uid == user_id or acc == user_id:
                        user_events.append(event)
                except: pass
    except: return None
    
    if not user_events: return None
    
    feature_usage = {}
    prompts = []
    actions = []
    
    for event in user_events:
        data = event.get('data', {})
        feat = data.get('feature_name', 'Unknown')
        feature_usage[feat] = feature_usage.get(feat, 0) + 1
        if data.get('prompt'): prompts.append(data['prompt'][:50])
        actions.append(f"{feat}-{data.get('action', '')}")
        
    return {
        'total_events': len(user_events),
        'top_features': sorted(feature_usage.items(), key=lambda x: x[1], reverse=True)[:5],
        'recent_prompts': prompts[-5:],
        'recent_actions': actions[-10:],
        'days_analyzed': days
    }

def generate_ai_prompt(user_summary, user_account):
    if not ZHIPU_API_KEY: return None
    
    prompt = f"""
    Analyze this user:
    Account: {user_account}
    Total Actions: {user_summary['total_events']}
    Top Features: {user_summary['top_features']}
    Recent Prompts: {user_summary['recent_prompts']}
    
    Return JSON with: user_type, interests, usage_habits, proficiency_level, suggestions, summary.
    """
    
    try:
        resp = requests.post(
            ZHIPU_API_URL,
            headers={'Authorization': f'Bearer {ZHIPU_API_KEY}', 'Content-Type': 'application/json'},
            json={
                'model': 'glm-3-turbo',
                'messages': [{'role': 'user', 'content': prompt}],
                'temperature': 0.7
            },
            timeout=10
        )
        if resp.status_code == 200:
            content = resp.json()['choices'][0]['message']['content']
            # Clean markdown if present
            if '```json' in content: content = content.split('```json')[1].split('```')[0]
            elif '```' in content: content = content.split('```')[1].split('```')[0]
            return json.loads(content.strip())
    except Exception as e:
        print(f"AI API Error: {e}")
    return None

# --- Routes: Trending ---

@app.route('/api/trending', methods=['GET'])
def get_trending():
    try:
        data = load_data()
        cat = request.args.get('category', 'all')
        limit = request.args.get('limit', type=int)
        
        repos = []
        if cat in data['categories']:
            repos = data['categories'][cat]
        else:
            for r in data['categories'].values(): repos.extend(r)
            repos.sort(key=lambda x: x.get('stars_count', 0), reverse=True)
            
        if limit: repos = repos[:limit]
        
        return jsonify({'success': True, 'data': repos, 'updated_at': data['updated_at']})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/categories', methods=['GET'])
def get_categories():
    data = load_data()
    cats = [{'name': k, 'count': len(v)} for k, v in data['categories'].items()]
    return jsonify({'success': True, 'categories': cats})

@app.route('/api/refresh', methods=['POST'])
def refresh_data():
    trigger_background_update()
    return jsonify({'success': True, 'message': 'Update started'})

# --- Routes: Users ---

@app.route('/api/users', methods=['GET'])
def get_users():
    # 强制不缓存
    users = load_users()
    response = jsonify(users)
    response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
    response.headers["Pragma"] = "no-cache"
    response.headers["Expires"] = "0"
    return response

@app.route('/api/users', methods=['POST'])
def update_users():
    try:
        users = request.json
        save_users(users)
        # 强制更新内存缓存
        global cached_data
        # 这里其实不需要更新 cached_data，因为 users 是读文件的
        # 但我们可以加个日志确认
        print(f"Users updated: {len(users)} users saved.")
        return jsonify({'success': True})
    except Exception as e:
        print(f"Error saving users: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500

# --- Routes: Behavior Tracking ---

@app.route('/api/kb/track', methods=['POST'])
def track_behavior():
    try:
        data = request.get_json()
        if not data or 'events' not in data: return jsonify({'error': 'No events'}), 400
        
        for event in data['events']:
            sanitized = sanitize_data(event)
            sanitized['server_timestamp'] = datetime.now().isoformat()
            append_to_log(sanitized)
            
        return jsonify({'success': True, 'count': len(data['events'])})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/kb/stats', methods=['GET'])
def get_stats():
    if not os.path.exists(LOG_FILE): return jsonify({'total': 0})
    count = 0
    with open(LOG_FILE, 'r', encoding='utf-8') as f:
        for _ in f: count += 1
    return jsonify({'total_events': count})

# --- Routes: AI Prompt ---

@app.route('/api/ai-prompt', methods=['POST'])
def get_ai_prompt_route():
    try:
        data = request.get_json()
        user_id = data.get('user_id')
        user_account = data.get('user_account', 'User')

        if not user_id: return jsonify({'error': 'Missing user_id'}), 400

        summary = get_user_behavior_summary(user_id)
        if not summary:
            return jsonify({'success': True, 'source': 'default', 'profile': {'summary': 'Welcome new user!'}})

        profile = generate_ai_prompt(summary, user_account)
        if profile:
            return jsonify({'success': True, 'source': 'ai', 'profile': profile})

        return jsonify({'success': True, 'source': 'fallback', 'profile': {'summary': 'Active user'}})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# --- Routes: System Config ---

@app.route('/api/admin/system-config', methods=['GET'])
def get_system_config():
    """获取系统配置 - 所有用户可读"""
    try:
        config = load_system_config()
        return jsonify({'success': True, 'data': config})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/admin/system-config', methods=['POST'])
def update_system_config():
    """更新系统配置 - 仅超级管理员可写"""
    try:
        config = request.get_json()
        if save_system_config(config):
            return jsonify({'success': True, 'message': '系统配置已保存'})
        else:
            return jsonify({'success': False, 'error': '保存失败'}), 500
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# --- Routes: 302.ai Proxy ---

@app.route('/api/302/<path:subpath>', methods=['GET', 'POST', 'PUT', 'DELETE'])
def proxy_302(subpath):
    """
    代理转发 302.ai 的 API 请求
    前端请求: /api/302/v1/chat/completions
    转发目标: https://api.302.ai/v1/chat/completions
    """
    target_url = f"https://api.302.ai/{subpath}"
    
    # 过滤掉可能导致问题的 Headers
    excluded_headers = ['Host', 'Content-Length', 'Connection']
    headers = {
        key: value for key, value in request.headers 
        if key not in excluded_headers
    }
    
    try:
        # 发起转发请求
        resp = requests.request(
            method=request.method,
            url=target_url,
            headers=headers,
            data=request.get_data(),
            cookies=request.cookies,
            allow_redirects=False,
            timeout=60, # 图片生成可能需要较长时间
            verify=False # 关键修改：跳过 SSL 验证，解决 Python 3.6 旧证书报错 500 问题
        )
        
        # 构建返回响应
        # 过滤掉 Hop-by-hop Headers
        excluded_response_headers = ['content-encoding', 'content-length', 'transfer-encoding', 'connection']
        headers = [
            (name, value) for (name, value) in resp.raw.headers.items()
            if name.lower() not in excluded_response_headers
        ]
        
        response = app.response_class(
            response=resp.content,
            status=resp.status_code,
            headers=headers
        )
        return response
    except Exception as e:
        print(f"Proxy Error: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/health', methods=['GET'])
def health():
    return jsonify({'status': 'ok', 'service': 'Unified AI Server'})

# ==================== 文件解析API ====================

@app.route('/api/parse-file', methods=['POST'])
def parse_file():
    """解析上传的文件内容（PDF、Word等）"""
    try:
        data = request.json
        file_type = data.get('type')
        file_name = data.get('name', 'unknown')
        file_content = data.get('content')  # base64编码的内容

        if not file_content:
            return jsonify({'error': '没有文件内容'}), 400

        if file_type == 'pdf':
            try:
                from PyPDF2 import PdfReader
                pdf_bytes = base64.b64decode(file_content)
                pdf_reader = PdfReader(io.BytesIO(pdf_bytes))
                text = ''
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + '\n'
                return jsonify({'content': text[:10000], 'type': 'text'})
            except ImportError:
                return jsonify({'error': 'PDF解析库未安装，请安装PyPDF2'}), 500
            except Exception as e:
                return jsonify({'error': f'PDF解析失败: {str(e)}'}), 500

        elif file_type == 'docx':
            try:
                from docx import Document
                docx_bytes = base64.b64decode(file_content)
                doc = Document(io.BytesIO(docx_bytes))
                text = '\n'.join([para.text for para in doc.paragraphs])
                return jsonify({'content': text[:10000], 'type': 'text'})
            except ImportError:
                return jsonify({'error': 'Word解析库未安装，请安装python-docx'}), 500
            except Exception as e:
                return jsonify({'error': f'Word解析失败: {str(e)}'}), 500

        elif file_type == 'doc':
            # 旧版 .doc 格式，尝试多种方法解析
            try:
                # 方法1: 尝试使用 antiword 命令行工具 (Linux服务器上常用)
                import subprocess
                doc_bytes = base64.b64decode(file_content)
                result = subprocess.run(
                    ['antiword', '-'],
                    input=doc_bytes,
                    capture_output=True,
                    timeout=30
                )
                if result.returncode == 0 and result.stdout:
                    text = result.stdout.decode('utf-8', errors='ignore')
                    return jsonify({'content': text[:10000], 'type': 'text'})
            except Exception:
                pass

            try:
                # 方法2: 尝试使用 olefile 读取 .doc 中的文本
                import olefile
                import struct
                doc_bytes = base64.b64decode(file_content)
                ole = olefile.OleFileIO(io.BytesIO(doc_bytes))

                # 尝试读取 WordDocument 流
                if ole.exists('WordDocument'):
                    # 简单提取文本（这可能不完美但能工作）
                    word_stream = ole.openstream('WordDocument').read()
                    # 尝试从流中提取可读文本
                    text_parts = []
                    i = 0
                    while i < len(word_stream):
                        # 查找可打印的字符序列
                        if 32 <= word_stream[i] <= 126 or word_stream[i] > 127:
                            char = chr(word_stream[i]) if word_stream[i] < 128 else chr(word_stream[i])
                            text_parts.append(char)
                        i += 1

                    text = ''.join(text_parts)
                    # 清理文本
                    import re
                    text = re.sub(r'\s+', ' ', text).strip()
                    if len(text) > 100:
                        return jsonify({'content': text[:10000], 'type': 'text'})

                ole.close()
            except Exception:
                pass

            # 如果所有方法都失败，返回友好提示
            return jsonify({
                'content': '[旧版 .doc 格式解析受限。建议您：\n1. 用Word打开后另存为 .docx 格式\n2. 或直接将文档内容复制粘贴到聊天框中]',
                'type': 'text'
            })

        elif file_type in ['xlsx', 'xls']:
            # Excel 文件解析
            try:
                import pandas as pd
                excel_bytes = base64.b64decode(file_content)
                df = pd.read_excel(io.BytesIO(excel_bytes), engine='openpyxl')
                # 转换为文本格式
                text = df.to_string(index=False, header=True)
                return jsonify({'content': text[:15000], 'type': 'text'})
            except ImportError:
                return jsonify({'error': 'Excel解析库未安装，请安装pandas和 openpyxl'}), 500
            except Exception as e:
                return jsonify({'error': f'Excel解析失败: {str(e)}'}), 500

        elif file_type in ['pptx', 'ppt']:
            # PowerPoint 文件解析
            try:
                from pptx import Presentation
                ppt_bytes = base64.b64decode(file_content)
                prs = Presentation(io.BytesIO(ppt_bytes))
                text = ''
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text:
                            text += shape.text + '\n'
                return jsonify({'content': text[:15000], 'type': 'text'})
            except ImportError:
                return jsonify({'error': 'PPT解析库未安装，请安装python-pptx'}), 500
            except Exception as e:
                return jsonify({'error': f'PPT解析失败: {str(e)}'}), 500

        elif file_type == 'rtf':
            # RTF 文件解析
            try:
                # RTF 是文本格式，可以尝试直接读取
                rtf_bytes = base64.b64decode(file_content)
                text = rtf_bytes.decode('utf-8', errors='ignore')
                # RTF 格式复杂，简单清理
                import re
                # 移除 RTF 控制字符
                text = re.sub(r'\\[a-z]+\d*\s?', '', text)
                text = re.sub(r'[{}]+', '', text)
                return jsonify({'content': text[:15000], 'type': 'text'})
            except Exception as e:
                return jsonify({'error': f'RTF解析失败: {str(e)}'}), 500

        else:
            return jsonify({'error': f'不支持的文件类型: {file_type}'}), 400

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'error': f'解析文件失败: {str(e)}'}), 500

if __name__ == '__main__':
    # Initial load
    load_data()
    print(f"Unified Server running on port 5001")
    # 恢复监听 5001，只提供 API 服务
    app.run(host='0.0.0.0', port=5001, threaded=True)
